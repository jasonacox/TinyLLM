"""
Document Generation Module for TinyLLM Chatbot

This module handles document generation functionality including PDF, Word, Excel, 
and PowerPoint document creation based on LLM responses.

Author: Jason A. Cox
3 Aug 2025
github.com/jasonacox/TinyLLM
"""

import os
import datetime
import uuid
from typing import Optional, Dict, Any

from app.core.config import log, debug
from app.core.llm import ask_llm

# Create directory for generated documents
# Get the chatbot directory (3 levels up from this file)
_current_file = os.path.abspath(__file__)
_chatbot_dir = os.path.dirname(os.path.dirname(os.path.dirname(_current_file)))
GENERATED_DOCS_DIR = os.path.join(_chatbot_dir, "generated_docs")
os.makedirs(GENERATED_DOCS_DIR, exist_ok=True)


class DocumentGenerator:
    """Handles document generation for various formats."""
    
    def __init__(self):
        # Map format aliases to canonical output types
        # Canonical types correspond to the actual file extensions produced
        self.format_aliases = {
            # PDF
            "pdf": "pdf",
            # Word / Docx
            "word": "docx",
            "doc": "docx",
            "docx": "docx",
            # Excel / Spreadsheet
            "excel": "xlsx",
            "spreadsheet": "xlsx",
            "xlsx": "xlsx",
            # PowerPoint
            "powerpoint": "pptx",
            "ppt": "pptx",
            "pptx": "pptx",
        }

    def canonical_format(self, doc_type: str) -> str:
        """Return the canonical format for a given document type alias.

        Falls back to 'pdf' if unknown.
        """
        if not doc_type:
            return "pdf"
        return self.format_aliases.get(doc_type.lower().strip(), "pdf")
    
    async def generate_document(self, session_id: str, llm_response: str, doc_request: Dict[str, Any], model: str) -> Optional[str]:
        """
        Generate document based on LLM response and user request.
        
        Args:
            session_id: Session identifier
            llm_response: The LLM's response content
            doc_request: Document request configuration
            model: LLM model to use for structuring
            
        Returns:
            str: Generated file path or None if failed
        """
        try:
            debug(f"Starting document generation for session {session_id[:8]}")
            debug(f"Document type: {doc_request['type']}")
            canonical = self.canonical_format(doc_request['type'])
            debug(f"Canonical document type: {canonical}")
            debug(f"Original prompt: {doc_request.get('original_prompt', 'Not provided')}")
            debug(f"LLM response length: {len(llm_response)} characters")
            debug(f"Model: {model}")
            
            log(f"Generating {doc_request['type'].upper()} document for session {session_id[:8]}")
            
            # Step 1: Ask LLM to structure the content for document
            debug("Step 1: Structuring content with LLM")
            # Provide a PowerPoint-specific outline format when generating PPTX, else use a general structure prompt
            if canonical == "pptx":
                structure_prompt = f"""
                You are drafting a presentation OUTLINE in Markdown for a PowerPoint deck.
                Follow these strict rules and output ONLY Markdown (no preamble, no commentary):
                - Start with ONE top-level H1 line: `# Presentation: <Short Title>`
                - Then create slides as H2 sections using: `## Slide N: <Slide Title>` (number consecutively)
                - Under each H2, include concise bullets using `- ` (3–6 bullets per slide). Keep bullets short (4–10 words).
                - If steps are required, you may use a numbered list `1.`, `2.`, etc.
                - For tabular data, use a Markdown table with a header row and a separator row of dashes.
                - For code, use fenced code blocks with triple backticks. Avoid images; if an image would help, use a bullet instead.
                - Do not include speaker notes, footers, or any text outside the outline. Avoid nested headings deeper than H2.
                - Target 5–10 slides unless the topic clearly needs fewer or more.

                Example format:
                # Presentation: Solar Energy 101
                ## Slide 1: What Is Solar?
                - Definition of solar power
                - How photovoltaic (PV) cells work
                - Key terms: PV, watt, kilowatt-hour

                ## Slide 2: Market Snapshot (2024)
                | Metric | Value | YoY |
                | --- | ---: | ---: |
                | Global capacity (GW) | 500 | +12% |
                | Avg. module price ($/W) | 0.18 | -15% |

                ## Slide 3: Sample Command
                ```
                pip install solar
                solar simulate --hours 24
                ```

                End of example.

                The user's original request: {doc_request["original_prompt"]}

                Here is draft content to incorporate faithfully (edit for slide-suitability and succinct bullets):
                {llm_response}

                Produce the final Markdown outline now, following the exact format above.
                """
            else:
                structure_prompt = f"""
                The user originally requested: {doc_request["original_prompt"]}
                
                Here is the response content: {llm_response}
                
                Please structure this content appropriately for a {doc_request["type"]} document.
                For the format, provide:
                - A clear title
                - Organized sections with headers
                - Any data that should be formatted as tables
                - Proper formatting suggestions
                
                Return the structured content in a clear, organized format.
                """
            
            debug(f"Sending structure prompt to LLM (length: {len(structure_prompt)} chars)")
            structured_content = await ask_llm(structure_prompt, model=model)
            debug(f"Received structured content from LLM (length: {len(structured_content)} chars)")
            
            # Step 2: Generate the actual document
            debug("Step 2: Creating document file")
            file_path = await self.create_document_file(session_id, structured_content, canonical)
            
            if file_path:
                debug(f"Document generation successful: {file_path}")
                log(f"Document generated successfully: {file_path}")
            else:
                debug("Document generation failed: file_path is None")
                log(f"Document generation failed for session {session_id[:8]}")
            
            return file_path
            
        except Exception as e:
            debug(f"Exception in generate_document: {type(e).__name__}: {e}")
            log(f"Document generation error: {e}")
            return None
    
    async def create_document_file(self, session_id: str, content: str, doc_type: str) -> Optional[str]:
        """Create the actual document file."""
        try:
            debug(f"Creating document file for session {session_id[:8]}")
            debug(f"Document type: {doc_type}")
            debug(f"Content length: {len(content)} characters")
            
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            # Use an opaque ID to avoid leaking session identifiers in filenames
            opaque_id = uuid.uuid4().hex[:8]
            filename = f"document_{opaque_id}_{timestamp}"
            debug(f"Generated filename: {filename}")
            
            # Check if output directory exists
            debug(f"Checking if output directory exists: {GENERATED_DOCS_DIR}")
            if not os.path.exists(GENERATED_DOCS_DIR):
                debug(f"Creating output directory: {GENERATED_DOCS_DIR}")
                os.makedirs(GENERATED_DOCS_DIR, exist_ok=True)
            
            # Branch based on canonical type
            if doc_type == "pdf":
                debug("Creating PDF document")
                file_path = await self.create_pdf_document(content, filename)
            elif doc_type == "docx":
                debug("Creating Word document")
                file_path = await self.create_word_document(content, filename)
            elif doc_type == "xlsx":
                debug("Creating Excel document")
                file_path = await self.create_excel_document(content, filename)
            elif doc_type == "pptx":
                debug("Creating PowerPoint document")
                file_path = await self.create_powerpoint_document(content, filename)
            else:
                debug(f"Unknown document type '{doc_type}', defaulting to PDF")
                file_path = await self.create_pdf_document(content, filename)
            
            if file_path:
                debug(f"Document file created successfully: {file_path}")
            else:
                debug("Document file creation failed")
            
            return file_path
            
        except Exception as e:
            debug(f"Exception in create_document_file: {type(e).__name__}: {e}")
            log(f"Document creation error: {e}")
            return None
    
    async def create_pdf_document(self, content: str, filename: str) -> Optional[str]:
        """Create PDF using reportlab."""
        try:
            debug(f"Starting PDF creation for filename: {filename}")
            debug(f"Content length: {len(content)} characters")
            
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, ListFlowable, ListItem, XPreformatted
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.lib import colors
            debug("Successfully imported reportlab modules")
            
            file_path = os.path.join(GENERATED_DOCS_DIR, f"{filename}.pdf")
            debug(f"Target file path: {file_path}")
            
            # Create PDF document
            debug("Creating SimpleDocTemplate")
            doc = SimpleDocTemplate(file_path, pagesize=letter)
            styles = getSampleStyleSheet()
            # Define small header for code and a code block style
            code_header_style = ParagraphStyle(
                name='CodeHeader', parent=styles['Normal'], fontSize=9, leading=11,
                textColor=colors.HexColor('#444444'), spaceAfter=4)
            code_block_style = ParagraphStyle(
                name='CodeBlock', parent=(styles['Code'] if 'Code' in styles else styles['Normal']),
                fontName='Courier', fontSize=9, leading=11)
            story = []

            # Helpers: lightweight Markdown handling
            import re

            def escape_text(s: str) -> str:
                # Minimal escaping for Paragraph XML
                return (
                    s.replace('&', '&amp;')
                     .replace('<', '&lt;')
                     .replace('>', '&gt;')
                )

            def inline_md_to_para(s: str) -> str:
                # Convert a small subset of Markdown into ReportLab Paragraph XML
                # 1) Escape
                s = escape_text(s)
                # 2) Bold/Italic
                s = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", s)
                s = re.sub(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", r"<i>\1</i>", s)
                # 3) Links: strip internal anchors to plain text; keep http(s) links clickable
                def _link_repl(m):
                    text = m.group(1)
                    target = m.group(2)
                    if target.startswith('#'):
                        return text  # drop anchor-only link
                    return f'<link href="{target}">{text}</link>'
                s = re.sub(r"\[([^\]]+)\]\((#[^)]+|https?://[^)]+)\)", _link_repl, s)
                return s

            # No anchor slugging needed when anchors are removed

            def is_md_header_sep(line: str) -> bool:
                # Detect Markdown table header separator: | --- | :---: | ---: |
                t = line.strip()
                if '|' not in t:
                    return False
                cells = [c.strip() for c in t.strip('|').split('|')]
                if not cells:
                    return False
                for c in cells:
                    if not re.fullmatch(r":?-{3,}:?", c):
                        return False
                return True

            def is_md_piped_row(line: str) -> bool:
                t = line.strip()
                if t.startswith('#'):
                    return False
                if t.count('|') < 2:
                    return False
                return t.startswith('|') or t.endswith('|')

            # Parse content into Flowables with Markdown support (headings, lists, tables)
            lines = content.split('\n')
            debug(f"Parsed content into {len(lines)} lines")
            i = 0
            processed_lines = 0
            # Simplified: no anchor tracking
            while i < len(lines):
                line = lines[i].strip()
                if not line:
                    story.append(Spacer(1, 0.15 * inch))
                    i += 1
                    continue
                # Fenced code block: ```lang (optional) ... ```
                if line.startswith("```"):
                    # extract language from fence if present
                    lang = None
                    import re as _re
                    m = _re.match(r"^```\s*([A-Za-z0-9_+\-\.]+)?", line)
                    if m:
                        lang = (m.group(1) or '').strip() or None
                    i += 1
                    code_lines = []
                    # accumulate until closing fence or EOF
                    while i < len(lines) and not lines[i].strip().startswith("```"):
                        code_lines.append(lines[i].rstrip('\n'))
                        i += 1
                    # skip closing fence if present
                    if i < len(lines) and lines[i].strip().startswith("```"):
                        i += 1
                    # Build header and shaded code block
                    header_text = f"<b>Code{f' ({lang})' if lang else ''}</b>"
                    story.append(Paragraph(header_text, code_header_style))
                    code_text = "\n".join(code_lines)
                    # Escape to avoid XML interpretation in XPreformatted
                    code_text_escaped = escape_text(code_text)
                    pre = XPreformatted(code_text_escaped, code_block_style)
                    tbl = Table([[pre]], colWidths=[doc.width])
                    tbl.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, -1), colors.whitesmoke),
                        ('BOX', (0, 0), (-1, -1), 0.5, colors.lightgrey),
                        ('LEFTPADDING', (0, 0), (-1, -1), 6),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                        ('TOPPADDING', (0, 0), (-1, -1), 6),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                    ]))
                    story.append(tbl)
                    story.append(Spacer(1, 0.12 * inch))
                    processed_lines += 1 + len(code_lines)
                    continue
                # Headings: '#', '##', or line fully wrapped in **...**
                if line.startswith('##'):
                    text = line.lstrip('#').strip()
                    story.append(Paragraph(inline_md_to_para(text), styles['Heading2']))
                    story.append(Spacer(1, 0.08 * inch))
                    i += 1
                    processed_lines += 1
                    continue
                if line.startswith('#'):
                    text = line.lstrip('#').strip()
                    story.append(Paragraph(inline_md_to_para(text), styles['Heading1']))
                    story.append(Spacer(1, 0.1 * inch))
                    i += 1
                    processed_lines += 1
                    continue
                # Bold-only line -> render as a normal paragraph with bold, not as a heading
                if line.startswith('**') and line.endswith('**') and len(line) > 4 and line.count('**') == 2:
                    text = line.strip('*').strip()
                    story.append(Paragraph(f"<b>{inline_md_to_para(text)}</b>", styles['Normal']))
                    story.append(Spacer(1, 0.06 * inch))
                    i += 1
                    processed_lines += 1
                    continue

                # Bullet lists: lines starting with *, -, or +; collect consecutive items
                if line.startswith(('* ', '- ', '+ ')):
                    bullets = []
                    while i < len(lines) and lines[i].strip().startswith(('* ', '- ', '+ ')):
                        bullets.append(lines[i].strip()[2:].strip())
                        i += 1
                    # Create a ListFlowable
                    items = [ListItem(Paragraph(inline_md_to_para(b), styles['Normal'])) for b in bullets]
                    story.append(ListFlowable(items, bulletType='bullet'))
                    story.append(Spacer(1, 0.08 * inch))
                    processed_lines += len(bullets)
                    continue

                # Markdown table: header row + separator row, followed by data rows
                if i + 1 < len(lines) and is_md_piped_row(line) and is_md_header_sep(lines[i + 1].strip()):
                    header_cells = [c.strip() for c in line.strip().strip('|').split('|')]
                    i += 2  # skip header and separator
                    data_rows = []
                    while i < len(lines) and is_md_piped_row(lines[i].strip()):
                        row_cells = [c.strip() for c in lines[i].strip().strip('|').split('|')]
                        # Convert row cells to Paragraphs to handle inline markdown (e.g., **bold**)
                        row_paras = [Paragraph(inline_md_to_para(c), styles['Normal']) for c in row_cells]
                        data_rows.append(row_paras)
                        i += 1
                    # Convert header cells to Paragraphs as well so **markers** don't appear
                    header_paras = [Paragraph(inline_md_to_para(c), styles['Normal']) for c in header_cells]
                    tbl_data = [header_paras] + data_rows
                    table = Table(tbl_data, hAlign='LEFT')
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 6),
                        ('TOPPADDING', (0, 1), (-1, -1), 4),
                        ('BOTTOMPADDING', (0, 1), (-1, -1), 4),
                    ]))
                    story.append(table)
                    story.append(Spacer(1, 0.12 * inch))
                    processed_lines += 2 + len(data_rows)
                    continue

                # Fallback: paragraph with inline markdown (anchor-only links already removed in inline_md_to_para)
                story.append(Paragraph(inline_md_to_para(lines[i].strip()), styles['Normal']))
                story.append(Spacer(1, 0.06 * inch))
                i += 1
                processed_lines += 1

            debug(f"Processed {processed_lines} content lines into {len(story)} story elements")
            debug("Building PDF document")
            doc.build(story)
            debug(f"PDF document built successfully: {filename}.pdf")
            
            return f"{filename}.pdf"
            
        except Exception as e:
            debug(f"Exception in create_pdf_document: {type(e).__name__}: {e}")
            log(f"PDF creation error: {e}")
            return None
    
    async def create_word_document(self, content: str, filename: str) -> Optional[str]:
        """Create a Word document rendering a useful subset of Markdown.
        Supports: headings (#, ##), bold (**), italic (*), bullet/numbered lists, and Markdown tables.
        """
        try:
            debug(f"Starting Word document creation for filename: {filename}")
            debug(f"Content length: {len(content)} characters")

            from docx import Document
            debug("Successfully imported python-docx")

            import re

            file_path = os.path.join(GENERATED_DOCS_DIR, f"{filename}.docx")
            debug(f"Target file path: {file_path}")

            # Create Word document
            debug("Creating new Word document")
            doc = Document()

            # Helpers ---------------------------------------------------------
            def add_markdown_runs(paragraph, text: str):
                """Add runs to a paragraph with basic Markdown emphasis parsing.
                Handles **bold** and *italic* (non-nested, best-effort).
                """
                if not text:
                    return
                # Split by strong/em emphasis tokens while preserving them
                # Example tokens: **bold**, *italic*, or plain text
                pattern = re.compile(r"(\*\*[^*]+\*\*|\*[^*]+\*)")
                pos = 0
                for m in pattern.finditer(text):
                    if m.start() > pos:
                        paragraph.add_run(text[pos:m.start()])
                    token = m.group(0)
                    if token.startswith("**") and token.endswith("**"):
                        run = paragraph.add_run(token[2:-2])
                        run.bold = True
                    elif token.startswith("*") and token.endswith("*"):
                        run = paragraph.add_run(token[1:-1])
                        run.italic = True
                    else:
                        paragraph.add_run(token)
                    pos = m.end()
                if pos < len(text):
                    paragraph.add_run(text[pos:])

            def is_table_header(line: str, next_line: Optional[str]) -> bool:
                if line is None or next_line is None:
                    return False
                if '|' not in line:
                    return False
                sep = next_line.strip()
                if not sep or '|' not in sep:
                    return False
                # Markdown table separator e.g., |---|:---:|---|
                return all(ch in "-|:+ " for ch in sep) and ('-' in sep)

            def split_pipes(row: str) -> list[str]:
                # keep empty leading/trailing cells
                parts = [p.strip() for p in row.split('|')]
                # If line starts/ends with pipe, first/last will be '' which we keep
                # Remove accidental double-pipe empties in middle only
                return [p for p in parts]

            # -----------------------------------------------------------------
            lines = content.split('\n')
            debug(f"Parsed content into {len(lines)} lines")

            i = 0
            processed_lines = 0
            while i < len(lines):
                raw = lines[i]
                line = raw.rstrip()
                if line.strip() == "":
                    # blank line - add spacing paragraph
                    doc.add_paragraph("")
                    i += 1
                    processed_lines += 1
                    continue

                # Headings: ###### to #
                m = re.match(r"^(#{1,6})\s+(.*)$", line)
                if m:
                    level = min(len(m.group(1)), 9)
                    text = m.group(2).strip()
                    doc.add_heading(text, level=level)
                    debug(f"Added heading level {level}: {text[:50]}...")
                    i += 1
                    processed_lines += 1
                    continue

                # Markdown table
                next_line = lines[i+1] if i + 1 < len(lines) else None
                if is_table_header(line, next_line):
                    header_cells = [c for c in split_pipes(line) if c != ""]
                    # advance past header and separator
                    i += 2
                    data_rows = []
                    while i < len(lines):
                        r = lines[i]
                        if '|' not in r or r.strip() == "":
                            break
                        cells = [c for c in split_pipes(r) if c != ""]
                        data_rows.append(cells)
                        i += 1
                    rows = [header_cells] + data_rows
                    if rows and header_cells:
                        cols = len(header_cells)
                        table = doc.add_table(rows=len(rows), cols=cols)
                        table.style = 'Table Grid'
                        # Fill header
                        for c_idx, txt in enumerate(header_cells):
                            p = table.cell(0, c_idx).paragraphs[0]
                            runp = p.add_run("")  # ensure first run exists
                            # Clear any default empty run text
                            if p.runs and p.runs[0].text == "":
                                p.runs[0].text = ""
                            add_markdown_runs(p, txt)
                            for run in p.runs:
                                run.bold = True
                        # Fill data rows
                        for r_idx, row_cells in enumerate(data_rows, start=1):
                            # pad row cells to header width to preserve blanks
                            for c_idx in range(cols):
                                cell_txt = row_cells[c_idx] if c_idx < len(row_cells) else ""
                                p = table.cell(r_idx, c_idx).paragraphs[0]
                                add_markdown_runs(p, cell_txt)
                    processed = 2 + len(data_rows)
                    processed_lines += processed
                    debug(f"Added Markdown table with {len(rows)} rows and {len(header_cells)} cols")
                    continue

                # Bullet list (- or *)
                if re.match(r"^\s*[-*]\s+", line):
                    while i < len(lines) and re.match(r"^\s*[-*]\s+", lines[i]):
                        item_text = re.sub(r"^\s*[-*]\s+", "", lines[i]).strip()
                        p = doc.add_paragraph(style='List Bullet')
                        add_markdown_runs(p, item_text)
                        i += 1
                        processed_lines += 1
                    continue

                # Numbered list (1. or 1)
                if re.match(r"^\s*\d+[\.)]\s+", line):
                    while i < len(lines) and re.match(r"^\s*\d+[\.)]\s+", lines[i]):
                        item_text = re.sub(r"^\s*\d+[\.)]\s+", "", lines[i]).strip()
                        p = doc.add_paragraph(style='List Number')
                        add_markdown_runs(p, item_text)
                        i += 1
                        processed_lines += 1
                    continue

                # Fallback: normal paragraph with inline emphasis
                p = doc.add_paragraph()
                add_markdown_runs(p, line.strip())
                i += 1
                processed_lines += 1

            debug(f"Processed {processed_lines} content lines")
            debug("Saving Word document")
            doc.save(file_path)
            debug(f"Word document saved successfully: {filename}.docx")

            return f"{filename}.docx"

        except Exception as e:
            debug(f"Exception in create_word_document: {type(e).__name__}: {e}")
            log(f"Word document creation error: {e}")
            return None
    
    async def create_excel_document(self, content: str, filename: str) -> Optional[str]:
        """Create Excel using openpyxl.""" 
        try:
            debug(f"Starting Excel document creation for filename: {filename}")
            debug(f"Content length: {len(content)} characters")
            
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
            debug("Successfully imported openpyxl")
            
            file_path = os.path.join(GENERATED_DOCS_DIR, f"{filename}.xlsx")
            debug(f"Target file path: {file_path}")
            
            # Create Excel workbook
            debug("Creating new Excel workbook")
            wb = Workbook()
            ws = wb.active
            ws.title = "Generated Content"
            debug("Created active worksheet: Generated Content")

            # Basic column sizing for readability
            ws.column_dimensions['A'].width = 50
            ws.column_dimensions['B'].width = 70
            ws.column_dimensions['C'].width = 30
            ws.column_dimensions['D'].width = 30
            
            # Parse content and add to spreadsheet with simple Markdown-table detection
            lines = content.split('\n')
            debug(f"Parsed content into {len(lines)} lines")
            row = 1

            import re  # local import for lightweight regex checks

            # Styles and helpers
            thin = Side(style='thin', color='999999')
            border_all = Border(left=thin, right=thin, top=thin, bottom=thin)
            header_fill = PatternFill("solid", fgColor="DDDDDD")

            def _clean_inline_md(text: str) -> str:
                # Remove inline emphasis markers and backticks for Excel plain text
                text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
                text = re.sub(r"\*(.*?)\*", r"\1", text)
                text = re.sub(r"`(.*?)`", r"\1", text)
                return text

            def _extract_entire_link(text: str):
                m = re.match(r"^\s*\[([^\]]+)\]\(([^)]+)\)\s*$", text)
                if m:
                    return m.group(1), m.group(2)
                return None

            def _is_entire_bold(text: str) -> bool:
                return bool(re.match(r"^\s*\*\*[^*].*\*\*\s*$", text))

            def _is_entire_italic(text: str) -> bool:
                return bool(re.match(r"^\s*\*[^*].*\*\s*$", text))

            def _is_piped_row(s: str) -> bool:
                # Row with at least two pipes and starts or ends with a pipe is considered tabular,
                # even if some cells are empty (e.g., "| Orange | | | |").
                if '|' not in s or s.startswith('#'):
                    return False
                if s.count('|') < 2:
                    return False
                t = s.strip()
                return t.startswith('|') or t.endswith('|')

            def _is_header_sep(s: str) -> bool:
                # Markdown header separator like: | --- | :---: | ---: |
                if '|' not in s:
                    return False
                segments = [seg.strip() for seg in s.strip().strip('|').split('|')]
                if len(segments) < 1:
                    return False
                for seg in segments:
                    if not seg:
                        return False
                    if not re.fullmatch(r":?-{3,}:?", seg):
                        return False
                return True

            i = 0
            processed_lines = 0
            in_code_block = False
            code_accum = []
            while i < len(lines):
                line_raw = lines[i]
                line = line_raw.strip()
                if not line:
                    i += 1
                    continue

                # Code block fences
                if line.startswith("```"):
                    if not in_code_block:
                        in_code_block = True
                        code_accum = []
                        i += 1
                        continue
                    else:
                        # closing fence: write accumulated code block
                        in_code_block = False
                        code_text = "\n".join(code_accum)
                        # Write code block merged across A:B for readability
                        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=2)
                        cell = ws.cell(row=row, column=1, value=code_text)
                        cell.font = Font(name='Consolas')
                        cell.alignment = Alignment(wrap_text=True, vertical='top')
                        cell.fill = PatternFill("solid", fgColor="F5F5F5")
                        cell.border = border_all
                        debug(f"Added code block at row {row} with {len(code_accum)} lines")
                        row += 1
                        i += 1
                        continue
                if in_code_block:
                    code_accum.append(line_raw.rstrip("\n"))
                    i += 1
                    continue

                # Try Markdown table detection first: header row + separator row
                if i + 1 < len(lines) and _is_piped_row(line) and _is_header_sep(lines[i + 1].strip()):
                    debug(f"Detected Markdown table starting at input line {i}")
                    # Parse header
                    header_cells = [c.strip() for c in line.strip().strip('|').split('|')]
                    sep_idx = i + 1
                    i = sep_idx + 1

                    # Collect data rows
                    data_rows = []
                    while i < len(lines) and _is_piped_row(lines[i].strip()):
                        cells = [c.strip() for c in lines[i].strip().strip('|').split('|')]
                        data_rows.append(cells)
                        i += 1

                    # Write header
                    for col, cell_value in enumerate(header_cells, 1):
                        cell = ws.cell(row=row, column=col, value=_clean_inline_md(cell_value))
                        cell.font = Font(bold=True)
                        cell.alignment = Alignment(vertical='center')
                        cell.fill = header_fill
                        cell.border = border_all
                    row += 1

                    # Write data rows, padding to header column count to preserve blank cells
                    header_cols = len(header_cells)
                    for r in data_rows:
                        for col in range(1, header_cols + 1):
                            cell_value = r[col-1] if col-1 < len(r) else ""
                            cell = ws.cell(row=row, column=col, value=_clean_inline_md(cell_value))
                            cell.alignment = Alignment(vertical='top', wrap_text=True)
                            cell.border = border_all
                        row += 1

                    processed_lines += 2 + len(data_rows)
                    debug(f"Wrote Markdown table: {1 + len(data_rows)} rows at sheet rows up to {row-1}")
                    continue

                # Fallback: treat as consecutive pipe-delimited rows block (must be 2+ rows)
                if _is_piped_row(line):
                    start = i
                    block = []
                    while i < len(lines) and _is_piped_row(lines[i].strip()):
                        block.append([c.strip() for c in lines[i].strip().strip('|').split('|')])
                        i += 1
                    if len(block) >= 2:
                        debug(f"Detected pipe table block from line {start} to {i-1} ({len(block)} rows)")
                        for r_cells in block:
                            for col, cell_value in enumerate(r_cells, 1):
                                ws.cell(row=row, column=col, value=cell_value)
                            row += 1
                        processed_lines += len(block)
                        continue
                    else:
                        # Not enough rows to be a table; treat the single line as normal text
                        i = start  # reset to original line index to handle as normal text below

                # Headings: #, ##, etc.
                m = re.match(r"^(#{1,6})\s+(.*)$", line)
                if m:
                    level = len(m.group(1))
                    text = _clean_inline_md(m.group(2).strip())
                    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=2)
                    cell = ws.cell(row=row, column=1, value=text)
                    size = 16 if level == 1 else 14 if level == 2 else 12
                    cell.font = Font(bold=True, size=size)
                    cell.alignment = Alignment(vertical='center')
                    debug(f"Added heading level {level} at row {row}: {text[:50]}...")
                    row += 1
                    processed_lines += 1
                    i += 1
                    continue

                # Bullet list (- or *) block
                if re.match(r"^\s*[-*]\s+", line):
                    while i < len(lines) and re.match(r"^\s*[-*]\s+", lines[i].strip()):
                        item_text = re.sub(r"^\s*[-*]\s+", "", lines[i].strip())
                        link = _extract_entire_link(item_text)
                        text_val = _clean_inline_md(item_text if not link else link[0])
                        cell = ws.cell(row=row, column=1, value=f"• {text_val}")
                        cell.alignment = Alignment(wrap_text=True)
                        if link:
                            cell.hyperlink = link[1]
                            cell.style = 'Hyperlink'
                        row += 1
                        i += 1
                        processed_lines += 1
                    continue

                # Numbered list block
                if re.match(r"^\s*\d+[\.)]\s+", line):
                    while i < len(lines) and re.match(r"^\s*\d+[\.)]\s+", lines[i].strip()):
                        m2 = re.match(r"^\s*(\d+)[\.)]\s+(.*)$", lines[i].strip())
                        num = m2.group(1)
                        item_text = m2.group(2)
                        link = _extract_entire_link(item_text)
                        text_val = _clean_inline_md(item_text if not link else link[0])
                        cell = ws.cell(row=row, column=1, value=f"{num}. {text_val}")
                        cell.alignment = Alignment(wrap_text=True)
                        if link:
                            cell.hyperlink = link[1]
                            cell.style = 'Hyperlink'
                        row += 1
                        i += 1
                        processed_lines += 1
                    continue

                # Entire-line hyperlink
                link = _extract_entire_link(line)
                if link:
                    cell = ws.cell(row=row, column=1, value=link[0])
                    cell.hyperlink = link[1]
                    cell.style = 'Hyperlink'
                    row += 1
                    processed_lines += 1
                    i += 1
                    continue

                # Entire-line bold/italic
                if _is_entire_bold(line):
                    text = _clean_inline_md(line)
                    cell = ws.cell(row=row, column=1, value=text)
                    cell.font = Font(bold=True)
                    cell.alignment = Alignment(wrap_text=True)
                    row += 1
                    processed_lines += 1
                    i += 1
                    continue
                if _is_entire_italic(line):
                    text = _clean_inline_md(line)
                    cell = ws.cell(row=row, column=1, value=text)
                    cell.font = Font(italic=True)
                    cell.alignment = Alignment(wrap_text=True)
                    row += 1
                    processed_lines += 1
                    i += 1
                    continue

                # Fallback: Regular paragraph in first column, cleaned
                cell = ws.cell(row=row, column=1, value=_clean_inline_md(line))
                cell.alignment = Alignment(wrap_text=True)
                debug(f"Added content at row {row}: {line[:50]}...")
                row += 1
                processed_lines += 1
                i += 1
            
            debug(f"Processed {processed_lines} content lines into {row-1} Excel rows")
            debug("Saving Excel document")
            wb.save(file_path)
            debug(f"Excel document saved successfully: {filename}.xlsx")
            
            return f"{filename}.xlsx"
            
        except Exception as e:
            debug(f"Exception in create_excel_document: {type(e).__name__}: {e}")
            log(f"Excel creation error: {e}")
            return None
    
    async def create_powerpoint_document(self, content: str, filename: str) -> Optional[str]:
        """Create PowerPoint using python-pptx with basic Markdown rendering.
        Supports: headings (#, ##), bullet/numbered lists, Markdown tables (preserves empty cells),
        fenced code blocks (monospaced shaded textbox), and inline bold/italic runs.
        """
        try:
            debug(f"Starting PowerPoint document creation for filename: {filename}")
            debug(f"Content length: {len(content)} characters")

            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.util import Inches
            debug("Successfully imported python-pptx")

            import re

            file_path = os.path.join(GENERATED_DOCS_DIR, f"{filename}.pptx")
            debug(f"Target file path: {file_path}")

            # Create PowerPoint presentation
            debug("Creating new PowerPoint presentation")
            prs = Presentation()

            # Helpers ---------------------------------------------------------
            def _disable_native_bullets(paragraph):
                """Disable native PowerPoint bullets on a paragraph to prevent double bullets."""
                try:
                    from pptx.oxml.ns import qn
                    from pptx.oxml.xmlchemy import OxmlElement
                    pPr = paragraph._p.get_or_add_pPr()
                    # Remove any existing bullet properties
                    for child in list(pPr):
                        if child.tag in {qn('a:buNone'), qn('a:buChar'), qn('a:buAutoNum'), qn('a:buBlip')}:
                            pPr.remove(child)
                    pPr.append(OxmlElement('a:buNone'))
                except Exception as _:
                    # If anything fails, ignore and continue (fallback to whatever default is)
                    pass
            def add_markdown_runs_ppt(paragraph, text: str):
                # Parse **bold** and *italic* into runs
                pattern = re.compile(r"(\*\*[^*]+\*\*|\*[^*]+\*)")
                pos = 0
                for m in pattern.finditer(text):
                    if m.start() > pos:
                        paragraph.add_run().text = text[pos:m.start()]
                    token = m.group(0)
                    if token.startswith("**") and token.endswith("**"):
                        r = paragraph.add_run()
                        r.text = token[2:-2]
                        r.font.bold = True
                    elif token.startswith("*") and token.endswith("*"):
                        r = paragraph.add_run()
                        r.text = token[1:-1]
                        r.font.italic = True
                    else:
                        paragraph.add_run().text = token
                    pos = m.end()
                if pos < len(text):
                    paragraph.add_run().text = text[pos:]

            def is_md_header_sep(line: str) -> bool:
                t = line.strip()
                if '|' not in t:
                    return False
                cells = [c.strip() for c in t.strip('|').split('|')]
                if not cells:
                    return False
                return all(re.fullmatch(r":?-{3,}:?", c) for c in cells)

            def is_md_piped_row(line: str) -> bool:
                t = line.strip()
                if t.startswith('#'):
                    return False
                if t.count('|') < 2:
                    return False
                return t.startswith('|') or t.endswith('|')

            def split_pipes(row: str) -> list[str]:
                return [p.strip() for p in row.split('|')]

            # -----------------------------------------------------------------
            lines = content.split('\n')
            debug(f"Parsed content into {len(lines)} lines")

            current_slide = None
            text_frame = None
            slide_count = 0
            processed_lines = 0
            i = 0
            in_code_block = False
            code_accum = []
            # Track next y for placing tables/code boxes on current slide
            slide_free_top = Inches(2.0)

            def ensure_slide(title_text: str = "Generated Content"):
                nonlocal current_slide, text_frame, slide_count, slide_free_top
                if current_slide is None:
                    slide_layout = prs.slide_layouts[1]  # Title and Content
                    current_slide = prs.slides.add_slide(slide_layout)
                    title = current_slide.shapes.title
                    title.text = title_text
                    tfph = current_slide.placeholders[1]
                    text_frame_local = tfph.text_frame
                    text_frame_local.clear()
                    # Ensure the initial paragraph has bullets disabled
                    _disable_native_bullets(text_frame_local.paragraphs[0])
                    slide_free_top = Inches(2.0)
                    text_frame = text_frame_local
                    slide_count += 1

            def add_title_slide(title_text: str = "Generated Content"):
                nonlocal current_slide, text_frame, slide_count, slide_free_top
                slide_layout = prs.slide_layouts[0]  # Title Slide
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title_text
                # Optional subtitle left blank; don't set text_frame so next content starts new slide
                current_slide = None
                text_frame = None
                slide_free_top = Inches(2.0)
                slide_count += 1

            while i < len(lines):
                line_raw = lines[i]
                line = line_raw.strip()
                if line == "":
                    i += 1
                    continue

                # Code fences
                if line.startswith("```"):
                    ensure_slide()
                    if not in_code_block:
                        in_code_block = True
                        code_accum = []
                        i += 1
                        continue
                    else:
                        # closing
                        in_code_block = False
                        code_text = "\n".join(code_accum)
                        # add textbox
                        width = Inches(8.0)
                        height = Inches(max(0.6, min(4.0, 0.3 * (len(code_accum) + 1))))
                        left = Inches(1.0)
                        top = slide_free_top
                        tb = current_slide.shapes.add_textbox(left, top, width, height)
                        tf = tb.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        r = p.add_run()
                        r.text = code_text
                        r.font.name = 'Consolas'
                        # light gray background by setting shape fill
                        fill = tb.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
                        slide_free_top = Inches(slide_free_top.inches + height.inches + 0.1)
                        i += 1
                        continue
                if in_code_block:
                    code_accum.append(line_raw.rstrip('\n'))
                    i += 1
                    continue

                # Headings (#..######) -> start new slide
                m = re.match(r"^(#{1,6})\s+(.*)$", line)
                if m:
                    # Start a new slide for any heading
                    title_text = m.group(2).strip()
                    # Normalize common prefixes from our outline prompt
                    # e.g., "Presentation: Title" -> "Title" and "Slide N: Title" -> "Title"
                    is_h1 = len(m.group(1)) == 1
                    raw_title = title_text
                    title_text = re.sub(r"^Presentation:\s*", "", title_text, flags=re.IGNORECASE)
                    title_text = re.sub(r"^Slide\s*\d+\s*:\s*", "", title_text, flags=re.IGNORECASE)
                    current_slide = None
                    if is_h1 and re.match(r"(?i)^Presentation:\s*", raw_title):
                        add_title_slide(title_text if title_text else "Presentation")
                    else:
                        ensure_slide(title_text if title_text else "Slide")
                    i += 1
                    processed_lines += 1
                    continue

                # Bold-only line (e.g., **Slide 1: Title**) -> start new slide
                if line.startswith('**') and line.endswith('**') and line.count('**') == 2 and len(line) > 4:
                    raw = line.strip('*').strip().strip('"')
                    if re.match(r"(?i)^Presentation:\s*", raw):
                        title_text = re.sub(r"^Presentation:\s*", "", raw, flags=re.IGNORECASE)
                        add_title_slide(title_text if title_text else "Presentation")
                    else:
                        title_text = re.sub(r"^Presentation:\s*", "", raw, flags=re.IGNORECASE)
                        title_text = re.sub(r"^Slide\s*\d+\s*:\s*", "", title_text, flags=re.IGNORECASE)
                        current_slide = None
                        ensure_slide(title_text if title_text else "Slide")
                    i += 1
                    processed_lines += 1
                    continue

                # Tables: header row + separator row
                next_line = lines[i+1] if i + 1 < len(lines) else None
                if next_line and is_md_piped_row(line) and is_md_header_sep(next_line.strip()):
                    ensure_slide()
                    header_cells = [c.strip() for c in line.strip().strip('|').split('|')]
                    i += 2
                    data_rows = []
                    while i < len(lines) and is_md_piped_row(lines[i].strip()):
                        cells = [c.strip() for c in lines[i].strip().strip('|').split('|')]
                        data_rows.append(cells)
                        i += 1
                    cols = len(header_cells)
                    rows = 1 + len(data_rows)
                    # place table
                    left = Inches(1.0)
                    top = slide_free_top
                    width = Inches(8.0)
                    height = Inches(max(1.0, min(5.0, 0.35 * rows)))
                    table_shape = current_slide.shapes.add_table(rows, cols, left, top, width, height)
                    table = table_shape.table
                    # header
                    for c_idx, txt in enumerate(header_cells):
                        cell = table.cell(0, c_idx)
                        cell.text = txt.replace('**', '')
                        for run in cell.text_frame.paragraphs[0].runs:
                            run.font.bold = True
                    # rows, padding to header width
                    for r_idx, row_cells in enumerate(data_rows, start=1):
                        for c_idx in range(cols):
                            val = row_cells[c_idx] if c_idx < len(row_cells) else ""
                            table.cell(r_idx, c_idx).text = val.replace('**', '')
                    slide_free_top = Inches(slide_free_top.inches + height.inches + 0.1)
                    processed_lines += 2 + len(data_rows)
                    continue

                # Bullet list block
                if re.match(r"^\s*[-*]\s+", line):
                    ensure_slide()
                    # Special-case: bullet that indicates a new slide, e.g., '* **Slide 2:** Amber'
                    # If detected, create a new slide titled from the trailing text after the bold 'Slide N:' label.
                    if re.match(r"^\s*[-*]\s+\*\*Slide[^*]*\*\*", line):
                        item_text = re.sub(r"^\s*[-*]\s+", "", line).strip()
                        mslide = re.match(r"^\*\*Slide[^*]*\*\*\s*:?:\s*(.*)$", item_text)
                        title_text = (mslide.group(1).strip().strip('"') if mslide else item_text)
                        title_text = re.sub(r"^Presentation:\s*", "", title_text, flags=re.IGNORECASE)
                        title_text = re.sub(r"^Slide\s*\d+\s*:\s*", "", title_text, flags=re.IGNORECASE)
                        current_slide = None
                        ensure_slide(title_text if title_text else "Slide")
                        i += 1
                        processed_lines += 1
                        continue
                    while i < len(lines) and re.match(r"^\s*[-*]\s+", lines[i].strip()):
                        item_text = re.sub(r"^\s*[-*]\s+", "", lines[i].strip())
                        # simple bullet by prefixing
                        p = text_frame.add_paragraph() if len(text_frame.paragraphs) else text_frame.paragraphs[0]
                        # clear any existing text
                        p.text = ""
                        _disable_native_bullets(p)
                        # add a bullet char and space, then the parsed runs
                        r = p.add_run()
                        r.text = "• "
                        add_markdown_runs_ppt(p, item_text)
                        i += 1
                        processed_lines += 1
                    continue

                # Numbered list block
                if re.match(r"^\s*\d+[\.)]\s+", line):
                    ensure_slide()
                    while i < len(lines) and re.match(r"^\s*\d+[\.)]\s+", lines[i].strip()):
                        m2 = re.match(r"^\s*(\d+)[\.)]\s+(.*)$", lines[i].strip())
                        num = m2.group(1)
                        item_text = m2.group(2)
                        p = text_frame.add_paragraph() if len(text_frame.paragraphs) else text_frame.paragraphs[0]
                        p.text = ""
                        _disable_native_bullets(p)
                        r = p.add_run()
                        r.text = f"{num}. "
                        add_markdown_runs_ppt(p, item_text)
                        i += 1
                        processed_lines += 1
                    continue

                # Fallback: normal paragraph with inline bold/italic
                ensure_slide()
                # First paragraph may exist and be empty; use it, else add new
                if len(text_frame.paragraphs) == 1 and text_frame.paragraphs[0].text == "":
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                _disable_native_bullets(p)
                # populate runs
                add_markdown_runs_ppt(p, line)
                processed_lines += 1
                i += 1

            if slide_count == 0:
                # Fallback to single slide
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = "Generated Content"
                slide.placeholders[1].text_frame.text = content
                slide_count = 1

            debug(f"Created {slide_count} slides from {processed_lines} content lines")
            debug("Saving PowerPoint presentation")
            prs.save(file_path)
            debug(f"PowerPoint presentation saved successfully: {filename}.pptx")

            return f"{filename}.pptx"

        except Exception as e:
            debug(f"Exception in create_powerpoint_document: {type(e).__name__}: {e}")
            log(f"PowerPoint creation error: {e}")
            return None
    
    def is_supported_format(self, doc_type: str) -> bool:
        """Check if document type is supported."""
        return doc_type.lower() in self.format_aliases
    
    def get_supported_formats(self) -> list:
        """Get list of supported document formats."""
        return list(self.format_aliases.keys())


# Global document generator instance
document_generator = DocumentGenerator()


async def generate_document_from_response(session_id: str, llm_response: str, doc_request: Dict[str, Any], model: str) -> Optional[str]:
    """
    Convenience function to generate document from LLM response.
    
    Args:
        session_id: Session identifier
        llm_response: The LLM's response content
        doc_request: Document request configuration
        model: LLM model to use
        
    Returns:
        str: Generated file path or None if failed
    """
    return await document_generator.generate_document(session_id, llm_response, doc_request, model)


def get_document_generator() -> DocumentGenerator:
    """Get the global document generator instance."""
    return document_generator
